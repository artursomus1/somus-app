"""
gerar_pptx_proposta.py
Gerador de apresentacoes PPTX para propostas de consorcio - Somus Capital.

Utiliza templates PPTX localizados em Corporate/Templates/ e substitui
os dados de demonstracao pelos valores reais do cliente, preservando
toda a formatacao original (fontes, cores, tamanhos, negrito, etc.).

Tipos de apresentacao:
    - Comparativa: 3 cenarios lado a lado (10 slides)
    - Proposta: 1 cenario detalhado (10-11 slides)

Subtipos disponiveis:
    CETHD, CETHD_Lance_Fixo, Lance_Fixo, PF_Tradicional, PJ_Tradicional
"""

import os
import re
import shutil
import logging
from datetime import datetime
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Configuracao
# ---------------------------------------------------------------------------

logger = logging.getLogger(__name__)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATES_DIR = os.path.join(BASE_DIR, "Corporate", "Templates")
DEFAULT_OUTPUT_DIR = os.path.join(BASE_DIR, "Corporate", "Propostas")

SUBTIPOS_VALIDOS = [
    "CETHD",
    "CETHD_Lance_Fixo",
    "Lance_Fixo",
    "PF_Tradicional",
    "PJ_Tradicional",
]

# PF_Tradicional nao possui o slide de Contexto na Proposta
SUBTIPOS_SEM_CONTEXTO = ["PF_Tradicional"]

# ---------------------------------------------------------------------------
# Formatacao brasileira
# ---------------------------------------------------------------------------


def fmt_brl(valor: float, com_centavos: bool = True) -> str:
    """Formata valor monetario no padrao brasileiro: R$ 1.234.567,89"""
    if com_centavos:
        inteiro = int(abs(valor))
        centavos = round((abs(valor) - inteiro) * 100)
        if centavos >= 100:
            inteiro += 1
            centavos = 0
        parte_inteira = f"{inteiro:,}".replace(",", ".")
        resultado = f"R$ {parte_inteira},{centavos:02d}"
    else:
        inteiro = int(round(abs(valor)))
        parte_inteira = f"{inteiro:,}".replace(",", ".")
        resultado = f"R$ {parte_inteira}"
    if valor < 0:
        resultado = f"-{resultado}"
    return resultado


def fmt_pct(valor: float, sufixo: str = "") -> str:
    """Formata percentual: 1,50% ou 19,56% a.a."""
    texto = f"{valor:.4f}".rstrip("0")
    if texto.endswith("."):
        texto += "0"
    partes = texto.split(".")
    if len(partes[1]) < 2:
        partes[1] = partes[1].ljust(2, "0")
    texto = ",".join(partes)
    resultado = f"{texto}%"
    if sufixo:
        resultado += f" {sufixo}"
    return resultado


def fmt_pct_mensal(valor: float) -> str:
    """Formata CET mensal: 0,6210% a.m."""
    return fmt_pct(valor, "a.m.")


def fmt_pct_anual(valor: float) -> str:
    """Formata CET anual: 7,71% a.a."""
    return fmt_pct(valor, "a.a.")


def fmt_pct_simples(valor: float) -> str:
    """Formata percentual simples sem sufixo: 19,00%"""
    partes = f"{valor:.2f}".split(".")
    return f"{partes[0]},{partes[1]}%"


def fmt_pct_inteiro(valor: float) -> str:
    """Formata percentual como inteiro se possivel: 25% ou 2,50%"""
    if valor == int(valor):
        return f"{int(valor)}%"
    return fmt_pct_simples(valor)


def fmt_mes_ordinal(mes: int) -> str:
    """Retorna '6 mes'."""
    return f"{mes}\u00ba m\u00eas"


# ---------------------------------------------------------------------------
# Substituicao de texto preservando formatacao
# ---------------------------------------------------------------------------


def _replace_in_paragraph(paragraph, old_text: str, new_text: str) -> bool:
    """Substitui old_text por new_text em um paragrafo, preservando formatacao."""
    full_text = paragraph.text
    if old_text not in full_text:
        return False
    for run in paragraph.runs:
        if old_text in run.text:
            run.text = run.text.replace(old_text, new_text)
            return True
    # Texto distribuido entre runs: consolidar no primeiro
    new_full = full_text.replace(old_text, new_text)
    if paragraph.runs:
        paragraph.runs[0].text = new_full
        for i in range(1, len(paragraph.runs)):
            paragraph.runs[i].text = ""
    return True


def _replace_in_shape(shape, old_text: str, new_text: str) -> bool:
    """Substitui old_text por new_text em um shape, preservando formatacao."""
    if not shape.has_text_frame:
        return False
    replaced = False
    for paragraph in shape.text_frame.paragraphs:
        if _replace_in_paragraph(paragraph, old_text, new_text):
            replaced = True
    return replaced


def _replace_full_shape_text(shape, new_text: str) -> bool:
    """Substitui todo o texto de um shape por new_text, preservando formatacao."""
    if not shape.has_text_frame:
        return False
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs:
            paragraph.runs[0].text = new_text
            for i in range(1, len(paragraph.runs)):
                paragraph.runs[i].text = ""
            return True
    return False


def _get_shape_text(shape) -> str:
    """Retorna o texto completo de um shape."""
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


# ---------------------------------------------------------------------------
# Regex patterns
# ---------------------------------------------------------------------------

RE_PARCELA_NUM = re.compile(r"Parcela\s+(\d+)")
RE_PARCELA_VALOR = re.compile(r"R\$\s*[\d\.]+,\d{2}")
RE_PCT_AM = re.compile(r"[\d,]+%\s*a\.m\.")
RE_PCT_AA = re.compile(r"[\d,]+%\s*a\.a\.")
RE_MES_ORDINAL = re.compile(r"(\d+)\u00ba?\s*m\u00eas", re.IGNORECASE)


def _is_month_year(text: str) -> bool:
    """Verifica se o texto parece ser um mes/ano como 'Marco 2026'."""
    meses = [
        "janeiro", "fevereiro", "mar\u00e7o", "marco", "abril", "maio", "junho",
        "julho", "agosto", "setembro", "outubro", "novembro", "dezembro",
    ]
    text_lower = text.strip().lower()
    return any(text_lower.startswith(mes) for mes in meses)


# ---------------------------------------------------------------------------
# Parcela replacement - by Parcela number matching
# ---------------------------------------------------------------------------


def _replace_parcela_values(slide, parcelas: List[float]):
    """
    Substitui valores de parcelas no slide.
    Constroi um mapa {parcela_num: value_shape} usando proximidade espacial:
    cada "Parcela N" shape tem um "R$ X.XXX,XX" shape na mesma linha (mesmo top).
    """
    parcela_num_shapes = {}   # {num: shape}
    value_shapes_by_top = {}  # {top: [shapes with R$ values]}

    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if not text:
            continue

        match = RE_PARCELA_NUM.match(text.strip())
        if match:
            num = int(match.group(1))
            parcela_num_shapes[num] = shape
        elif RE_PARCELA_VALOR.match(text.strip()):
            top = shape.top
            if top not in value_shapes_by_top:
                value_shapes_by_top[top] = []
            value_shapes_by_top[top].append(shape)

    # Match each "Parcela N" to its closest R$ value shape at the same top
    for num, num_shape in parcela_num_shapes.items():
        top = num_shape.top
        if top in value_shapes_by_top:
            # Find the value shape closest horizontally to the parcela label
            closest = min(
                value_shapes_by_top[top],
                key=lambda s: abs(s.left - num_shape.left),
            )
            idx = num - 1
            if 0 <= idx < len(parcelas):
                _replace_full_shape_text(closest, fmt_brl(parcelas[idx]))
            else:
                _replace_full_shape_text(closest, "")


# ---------------------------------------------------------------------------
# Operation slide: label-value by vertical proximity within same column
# ---------------------------------------------------------------------------


def _build_label_value_map(slide) -> Dict[str, object]:
    """
    Constroi um mapa de {label_key: value_shape} para slides de operacao.
    Labels e valores estao organizados em colunas (left column e right column).
    Cada label tem um valor logo abaixo (proximo top, mesmo left range).
    """
    shapes_with_text = []
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if text:
            shapes_with_text.append((shape, text))

    # Separar em colunas baseado na posicao horizontal
    # Left column: left < 5000000 (arbitrary midpoint)
    # Right column: left >= 5000000
    MID_X = 5_000_000

    left_col = [(s, t) for s, t in shapes_with_text if s.left < MID_X]
    right_col = [(s, t) for s, t in shapes_with_text if s.left >= MID_X]

    # Sort each column by top position
    left_col.sort(key=lambda x: x[0].top)
    right_col.sort(key=lambda x: x[0].top)

    result = {}

    # Process left column: label-value pairs
    _extract_label_values(left_col, result)

    # Process right column: label-value pairs
    _extract_label_values(right_col, result)

    return result


def _extract_label_values(col_shapes, result: dict):
    """
    Extrai pares label-valor de uma coluna de shapes ordenados por top.
    O padrao e: label shape seguido de value shape.
    """
    last_label = ""
    for shape, text in col_shapes:
        text_lower = text.lower()

        # Identify labels
        if "valor total da carta" in text_lower and not text.startswith("R$"):
            last_label = "valor_carta"
        elif "lance embutido" in text_lower and "%" in text and "R$" not in text:
            last_label = "lance_embutido"
            result["lance_embutido_label_shape"] = shape
        elif "cr\u00e9dito dispon\u00edvel" in text_lower and not text.startswith("R$"):
            last_label = "credito_disponivel"
        elif "lance livre" in text_lower and "%" in text and "R$" not in text:
            last_label = "lance_livre"
            result["lance_livre_label_shape"] = shape
        elif "lance com recurso" in text_lower and "%" in text and "R$" not in text:
            last_label = "lance_proprio"
            result["lance_proprio_label_shape"] = shape
        elif "alavancagem" in text_lower and not text.startswith("R$"):
            last_label = "alavancagem"
        elif "cr\u00e9dito efetivo" in text_lower and not text.startswith("R$"):
            last_label = "credito_efetivo"
        elif "taxa de administra" in text_lower:
            last_label = "taxa_adm"
        elif "fundo de reserva" in text_lower:
            last_label = "fundo_reserva"
        elif "prazo total" in text_lower:
            last_label = "prazo_total"
        elif "expectativa de contempla" in text_lower:
            last_label = "contemplacao"
        elif "custo efetivo total" in text_lower:
            last_label = "cet"
        elif last_label and _is_value_text(text):
            # This is the value for the last label
            result[last_label] = shape
            # Special: CET can be anual or mensal
            if last_label == "cet":
                if RE_PCT_AA.search(text):
                    result["cet_anual"] = shape
                elif RE_PCT_AM.search(text):
                    result["cet_mensal"] = shape
                last_label = "cet"  # Keep for next CET value
            else:
                last_label = ""
        elif last_label == "" and not _is_label_text(text):
            pass  # Skip description/sub-text shapes
        else:
            # Description text below a value - reset
            if last_label and not _is_value_text(text) and not _is_label_text(text):
                last_label = ""


def _is_value_text(text: str) -> bool:
    """Check if text looks like a value (R$, percentage, months, ordinal month)."""
    t = text.strip()
    if t.startswith("R$"):
        return True
    if re.match(r"^\d+[,\.]\d+%", t):
        return True
    if re.match(r"^\d+%$", t):
        return True
    if re.match(r"^\d+ meses?$", t):
        return True
    if RE_MES_ORDINAL.match(t):
        return True
    if t.startswith("M\u00eas "):
        return True
    return False


def _is_label_text(text: str) -> bool:
    """Check if text looks like a label."""
    labels = [
        "valor total", "lance embutido", "cr\u00e9dito dispon",
        "lance livre", "alavancagem", "cr\u00e9dito efetivo",
        "taxa de administra", "fundo de reserva", "prazo total",
        "expectativa", "custo efetivo", "lance com recurso",
        "valores em reais", "taxas", "prazos",
    ]
    tl = text.lower()
    return any(l in tl for l in labels)


def _apply_operacao_values(slide, cenario: dict, tipo: str = "comparativa"):
    """
    Substitui valores no slide de operacao usando o mapa label-value.
    """
    lv_map = _build_label_value_map(slide)

    # Valor total da carta
    if "valor_carta" in lv_map:
        _replace_full_shape_text(lv_map["valor_carta"], fmt_brl(cenario["valor_carta"], com_centavos=False))

    # Lance embutido valor
    if "lance_embutido" in lv_map:
        _replace_full_shape_text(lv_map["lance_embutido"], fmt_brl(cenario["lance_embutido_valor"], com_centavos=False))

    # Lance embutido label (update percentage)
    if "lance_embutido_label_shape" in lv_map:
        pct_str = fmt_pct_inteiro(cenario["lance_embutido_pct"])
        _replace_full_shape_text(lv_map["lance_embutido_label_shape"], f"Lance embutido ({pct_str})")

    # Credito disponivel
    if "credito_disponivel" in lv_map:
        _replace_full_shape_text(lv_map["credito_disponivel"], fmt_brl(cenario["credito_disponivel"], com_centavos=False))

    # Lance livre (Comparativa)
    if "lance_livre" in lv_map:
        _replace_full_shape_text(lv_map["lance_livre"], fmt_brl(cenario.get("lance_livre_valor", 0), com_centavos=False))
    if "lance_livre_label_shape" in lv_map:
        pct_str = fmt_pct_inteiro(cenario.get("lance_livre_pct", 0))
        _replace_full_shape_text(lv_map["lance_livre_label_shape"], f"Lance livre ({pct_str})")

    # Lance com recurso proprio (Proposta)
    if "lance_proprio" in lv_map:
        _replace_full_shape_text(lv_map["lance_proprio"], fmt_brl(cenario.get("lance_proprio_valor", 0), com_centavos=False))
    if "lance_proprio_label_shape" in lv_map:
        pct_str = fmt_pct_inteiro(cenario.get("lance_proprio_pct", 0))
        _replace_full_shape_text(lv_map["lance_proprio_label_shape"], f"Lance com recurso pr\u00f3prio ({pct_str})")

    # Alavancagem (Comparativa)
    if "alavancagem" in lv_map:
        _replace_full_shape_text(lv_map["alavancagem"], fmt_brl(cenario.get("alavancagem", 0), com_centavos=False))

    # Credito efetivo (Proposta)
    if "credito_efetivo" in lv_map:
        _replace_full_shape_text(lv_map["credito_efetivo"], fmt_brl(cenario.get("credito_efetivo", 0), com_centavos=False))

    # Taxa de administracao
    if "taxa_adm" in lv_map:
        _replace_full_shape_text(lv_map["taxa_adm"], fmt_pct_simples(cenario["taxa_adm"]))

    # Fundo de reserva
    if "fundo_reserva" in lv_map:
        _replace_full_shape_text(lv_map["fundo_reserva"], fmt_pct_simples(cenario["fundo_reserva"]))

    # Prazo total
    if "prazo_total" in lv_map:
        _replace_full_shape_text(lv_map["prazo_total"], f"{cenario['prazo_total']} meses")

    # Contemplacao
    if "contemplacao" in lv_map:
        text = _get_shape_text(lv_map["contemplacao"])
        if text.startswith("M\u00eas") or text.startswith("Mes"):
            _replace_full_shape_text(lv_map["contemplacao"], f"M\u00eas {cenario['contemplacao_mes']}")
        else:
            _replace_full_shape_text(lv_map["contemplacao"], fmt_mes_ordinal(cenario["contemplacao_mes"]))

    # CET anual
    if "cet_anual" in lv_map:
        _replace_full_shape_text(lv_map["cet_anual"], fmt_pct_anual(cenario["cet_anual"]))

    # CET mensal
    if "cet_mensal" in lv_map:
        _replace_full_shape_text(lv_map["cet_mensal"], fmt_pct_mensal(cenario["cet_mensal"]))


# ---------------------------------------------------------------------------
# Processamento - COMPARATIVA
# ---------------------------------------------------------------------------


def _process_comparativa_cover(slide, data: dict):
    """Slide 1: capa da Comparativa."""
    n_cenarios = len(data["cenarios"])
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if not text:
            continue
        if "cen\u00e1rio" in text.lower() or "cenario" in text.lower():
            _replace_full_shape_text(shape, f"{n_cenarios} cen\u00e1rios  \u2014  {data['cliente_nome']}")
        elif _is_month_year(text):
            _replace_full_shape_text(shape, data["data_mes_ano"])


def _process_comparativa_visoes(slide, data: dict):
    """Slide 2: comparativo das visoes (3 colunas)."""
    cenarios = data["cenarios"]

    # Agrupar shapes em 3 colunas pelo eixo X
    shapes_with_text = []
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if text:
            shapes_with_text.append((shape, text))

    # Encontrar titulos "Contemplacao no Xo mes" para delimitar colunas
    # Cada bloco comeca com esse titulo e contem os valores abaixo
    col_blocks = []  # list of (x_start, x_end, shapes)

    # Identificar 3 blocos por "Contemplacao no"
    contemp_shapes = [(s, t) for s, t in shapes_with_text if "Contempla\u00e7\u00e3o no" in t]
    contemp_shapes.sort(key=lambda x: x[0].left)

    if len(contemp_shapes) >= 3:
        # Usar posicoes X dos titulos para definir colunas
        col_x = [s.left for s, _ in contemp_shapes]
        # Assign each shape to the nearest column
        for ci, cenario in enumerate(cenarios):
            col_left = col_x[ci]
            # Tolerance: shapes within ~3000000 EMUs of column start
            col_right = col_x[ci + 1] if ci + 1 < len(col_x) else 99999999
            col_shapes = [
                (s, t) for s, t in shapes_with_text
                if col_left - 200000 <= s.left < col_right - 200000
            ]
            col_shapes.sort(key=lambda x: x[0].top)
            _replace_visoes_column(col_shapes, cenario)


def _replace_visoes_column(col_shapes, cenario: dict):
    """Substitui valores em uma coluna do slide de visoes."""
    seen_cet = 0
    for shape, text in col_shapes:
        # Titulo "Contemplacao no Xo mes"
        if "Contempla\u00e7\u00e3o no" in text:
            _replace_full_shape_text(shape, f"Contempla\u00e7\u00e3o no {cenario['contemplacao_mes']}\u00ba m\u00eas")
        # Mes ordinal isolado
        elif RE_MES_ORDINAL.match(text.strip()) and "Contempla" not in text:
            _replace_full_shape_text(shape, fmt_mes_ordinal(cenario["contemplacao_mes"]))
        # CET mensal
        elif RE_PCT_AM.search(text) and "CET" not in text:
            _replace_full_shape_text(shape, fmt_pct_mensal(cenario["cet_mensal"]))
        # CET anual
        elif RE_PCT_AA.search(text) and "CET" not in text:
            _replace_full_shape_text(shape, fmt_pct_anual(cenario["cet_anual"]))
        # Lance livre: "25%  --  R$ 375.000"
        elif "%" in text and "R$" in text and len(text) < 50:
            pct_str = fmt_pct_inteiro(cenario["lance_livre_pct"])
            _replace_full_shape_text(shape, f"{pct_str}  \u2014  {fmt_brl(cenario['lance_livre_valor'], com_centavos=False)}")
        # Credito disponivel (R$ sozinho, tipicamente o primeiro R$ puro)
        elif text.startswith("R$") and "Cr\u00e9dito" not in text and "Alavancagem" not in text:
            # Distinguir credito vs alavancagem pela ordem vertical
            # O primeiro R$ puro na coluna e credito, o segundo e alavancagem
            pass
        # Labels - skip
        elif any(lbl in text.lower() for lbl in ["contempla\u00e7\u00e3o estimada", "cet mensal", "cet anual", "lance livre", "cr\u00e9dito dispon", "alavancagem"]):
            pass

    # Segunda passagem: substituir valores R$ puros na ordem (credito, alavancagem)
    r_values = [(shape, text) for shape, text in col_shapes
                if text.startswith("R$") and "%" not in text and len(text) < 30]
    if len(r_values) >= 1:
        _replace_full_shape_text(r_values[0][0], fmt_brl(cenario["credito_disponivel"], com_centavos=False))
    if len(r_values) >= 2:
        _replace_full_shape_text(r_values[1][0], fmt_brl(cenario["alavancagem"], com_centavos=False))


def _process_comparativa_operacao(slide, cenario: dict, op_num: int):
    """Slides 3, 5, 7: detalhe da operacao."""
    # Update title
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if "Opera\u00e7\u00e3o" in text and "Contempla\u00e7\u00e3o" in text:
            _replace_full_shape_text(shape, f"Opera\u00e7\u00e3o {op_num}  \u2014  Contempla\u00e7\u00e3o no {cenario['contemplacao_mes']}\u00ba m\u00eas")
            break

    # Replace all values using label-value map
    _apply_operacao_values(slide, cenario, tipo="comparativa")


def _process_comparativa_parcelas(slide, cenario: dict, op_num: int):
    """Slides 4, 6, 8: tabela de parcelas."""
    # Update title
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if "Opera\u00e7\u00e3o" in text and "Tabela de parcelas" in text:
            _replace_full_shape_text(shape, f"Opera\u00e7\u00e3o {op_num}  \u2014  Contempla\u00e7\u00e3o no {cenario['contemplacao_mes']}\u00ba m\u00eas  \u2014  Tabela de parcelas  (1/1)")
        elif "Corre\u00e7\u00e3o anual" in text:
            corr_str = fmt_pct_inteiro(cenario["correcao_anual"])
            _replace_full_shape_text(shape, f"Corre\u00e7\u00e3o anual pr\u00e9-fixada de {corr_str}")

    # Replace parcela values
    _replace_parcela_values(slide, cenario["parcelas"])


# ---------------------------------------------------------------------------
# Processamento - PROPOSTA
# ---------------------------------------------------------------------------


def _process_proposta_cover(slide, data: dict):
    """Slide 1: capa da Proposta."""
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if not text:
            continue
        text_lower = text.lower()
        if _is_month_year(text):
            _replace_full_shape_text(shape, data["data_mes_ano"])
        elif ("cr\u00e9dito" not in text_lower and "credito" not in text_lower
              and "proposta" not in text_lower and "estruturado" not in text_lower
              and len(text) < 50 and not any(c.isdigit() for c in text)
              and "|" not in text):
            _replace_full_shape_text(shape, data["cliente_nome"])


def _process_proposta_contexto(slide, data: dict):
    """
    Slide 2: contexto da operacao (nao existe em PF_Tradicional).
    O slide tem DUAS colunas (esquerda: AMBIENTE, direita: OBJETIVO)
    e uma faixa inferior com metricas e cronograma.
    Processa cada secao independentemente usando o conteudo dos shapes.
    """
    cenario = data["cenario"]

    # Separar shapes em coluna esquerda, coluna direita, e metricas/bottom
    MID_X = 5_000_000  # Divisor horizontal entre colunas
    METRICS_TOP = 4_500_000  # Metricas e cronograma ficam na parte inferior

    all_shapes = [
        (s, _get_shape_text(s)) for s in slide.shapes
        if s.has_text_frame and _get_shape_text(s)
        and "SOMUS CAPITAL" not in _get_shape_text(s)
        and "somuscapital" not in _get_shape_text(s)
    ]

    left_col = sorted(
        [(s, t) for s, t in all_shapes if s.left < MID_X and s.top < METRICS_TOP],
        key=lambda x: x[0].top,
    )
    right_col = sorted(
        [(s, t) for s, t in all_shapes if s.left >= MID_X],
        key=lambda x: x[0].top,
    )
    bottom_shapes = sorted(
        [(s, t) for s, t in all_shapes if s.top >= METRICS_TOP and s.left < MID_X],
        key=lambda x: (x[0].top, x[0].left),
    )

    # --- Coluna esquerda: AMBIENTE DO NEGOCIO ---
    state = "init"
    for shape, text in left_col:
        if "Contexto da opera" in text or "AMBIENTE DO NEG" in text:
            state = "expect_empresa"
            continue
        if state == "expect_empresa":
            _replace_full_shape_text(shape, data.get("empresa_nome", text))
            state = "expect_local"
            continue
        if state == "expect_local":
            _replace_full_shape_text(shape, data.get("localizacao", text))
            state = "expect_desc"
            continue
        if state == "expect_desc":
            _replace_full_shape_text(shape, data.get("descricao_operacao", text))
            state = "done"
            continue

    # --- Metricas (bottom center) ---
    for shape, text in bottom_shapes:
        if text.startswith("R$"):
            _replace_full_shape_text(shape, fmt_brl(cenario["valor_carta"], com_centavos=False))
        elif text.isdigit():
            _replace_full_shape_text(shape, str(cenario["prazo_total"]))
        elif text.startswith("M\u00eas ") and len(text) < 15:
            _replace_full_shape_text(shape, f"M\u00eas {cenario['contemplacao_mes']}")

    # --- Coluna direita: OBJETIVO + CRONOGRAMA ---
    state = "init"
    for shape, text in right_col:
        if "OBJETIVO DA PROPOSTA" in text:
            state = "expect_obj_titulo"
            continue
        if state == "expect_obj_titulo":
            _replace_full_shape_text(shape, data.get("objetivo_titulo", text))
            state = "expect_obj_horizonte"
            continue
        if state == "expect_obj_horizonte":
            _replace_full_shape_text(shape, data.get("horizonte_texto", text))
            state = "expect_obj_desc"
            continue
        if state == "expect_obj_desc":
            _replace_full_shape_text(shape, data.get("objetivo_descricao", text))
            state = "after_objetivo"
            continue

        # Cronograma bullets
        if text in ("01", "02", "03"):
            continue

        if "Carta de cr\u00e9dito de R$" in text:
            carta_str = fmt_brl(cenario["valor_carta"], com_centavos=False)
            corr_str = fmt_pct_inteiro(cenario["correcao_anual"])
            _replace_full_shape_text(shape, f"Carta de cr\u00e9dito de {carta_str} com corre\u00e7\u00e3o anual de {corr_str} pr\u00e9-fixado")

        elif "Contempla\u00e7\u00e3o estimada no m\u00eas" in text:
            emb_str = fmt_pct_inteiro(cenario["lance_embutido_pct"])
            prop_str = fmt_pct_inteiro(cenario.get("lance_proprio_pct", 0))
            _replace_full_shape_text(shape, f"Contempla\u00e7\u00e3o estimada no m\u00eas {cenario['contemplacao_mes']} via lance embutido ({emb_str}) + lance pr\u00f3prio ({prop_str})")

        elif "Cr\u00e9dito efetivo de R$" in text:
            efetivo = cenario.get("credito_efetivo", 0)
            total_lances = cenario.get("total_lances", 0)
            _replace_full_shape_text(shape, f"Cr\u00e9dito efetivo de {fmt_brl(efetivo, com_centavos=False)} ap\u00f3s lances totais de {fmt_brl(total_lances, com_centavos=False)}")


def _process_proposta_dados(slide, data: dict):
    """Slide de 'Principais dados': usa o mesmo label-value map."""
    cenario = data["cenario"]
    _apply_operacao_values(slide, cenario, tipo="proposta")


def _find_value_below(label_shape, all_shapes):
    """
    Encontra o shape de valor logo abaixo de um label shape.
    O valor esta na mesma coluna (mesmo left) mas na proxima linha (top maior).
    """
    candidates = [
        s for s in all_shapes
        if s.has_text_frame and _get_shape_text(s)
        and s.top > label_shape.top
        and s.left == label_shape.left
    ]
    if candidates:
        return min(candidates, key=lambda s: s.top)
    return None


def _process_proposta_estrutura(slide, data: dict):
    """Slide de 'Estrutura da operacao'."""
    cenario = data["cenario"]

    all_shapes_list = list(slide.shapes)
    shapes_ordered = sorted(
        [s for s in all_shapes_list if s.has_text_frame and _get_shape_text(s)],
        key=lambda s: (s.top, s.left),
    )

    # --- Replace labels/values by content matching ---
    for shape in shapes_ordered:
        text = _get_shape_text(shape)
        if not text:
            continue

        if "SOMUS CAPITAL" in text or "somuscapital" in text:
            continue

        # Titulo do consorcio estruturado
        if "Cons\u00f3rcio estruturado" in text:
            adm = data.get("administradora", "")
            _replace_full_shape_text(shape, f"Cons\u00f3rcio estruturado {adm}")
            continue

        # Subtitulo com correcao
        if "Grupos pr\u00e9-fixados" in text:
            corr_str = fmt_pct_inteiro(cenario["correcao_anual"])
            _replace_full_shape_text(shape, f"Grupos pr\u00e9-fixados  \u2014  Corre\u00e7\u00e3o anual de {corr_str} pr\u00e9-fixado")
            continue

        # Descricao longa
        if "A opera\u00e7\u00e3o utiliza" in text:
            _replace_full_shape_text(shape, cenario.get("estrutura_descricao", text))
            continue

        # Label-value pairs: label on top, value below at same left
        if "Administradora" in text and len(text) < 20:
            val_shape = _find_value_below(shape, all_shapes_list)
            if val_shape:
                _replace_full_shape_text(val_shape, data.get("administradora", ""))
            continue

        if "Corre\u00e7\u00e3o anual" in text and "Grupos" not in text and len(text) < 20:
            val_shape = _find_value_below(shape, all_shapes_list)
            if val_shape:
                corr_str = fmt_pct_inteiro(cenario["correcao_anual"])
                _replace_full_shape_text(val_shape, f"{corr_str} pr\u00e9-fixado")
            continue

        if "Capta\u00e7\u00e3o em at\u00e9" in text:
            val_shape = _find_value_below(shape, all_shapes_list)
            if val_shape:
                meses = cenario.get("captacao_meses", cenario["contemplacao_mes"])
                _replace_full_shape_text(val_shape, f"{meses} meses")
            continue

        if "Total de lances" in text and len(text) < 20:
            val_shape = _find_value_below(shape, all_shapes_list)
            if val_shape:
                total = cenario.get("total_lances", 0)
                _replace_full_shape_text(val_shape, fmt_brl(total, com_centavos=False))
            continue

        # Cronograma resumo
        if "Contempla\u00e7\u00e3o no m\u00eas" in text and "via lance" in text.lower():
            _replace_full_shape_text(shape, cenario.get("cronograma_resumo", text))
            continue

        # Timeline step 1: "Meses 1 a 12"
        if re.match(r"Meses 1 \u00e0", text):
            _replace_full_shape_text(shape, f"Meses 1 \u00e0 {cenario['contemplacao_mes']}")
            continue

        # Pre-contemplacao descricao
        if text.startswith("Parcelas de R$") and "ac\u00famulo" in text:
            parcela_pre = cenario["parcelas"][0] if cenario["parcelas"] else 0
            emb_str = fmt_pct_inteiro(cenario["lance_embutido_pct"])
            _replace_full_shape_text(shape, f"Parcelas de {fmt_brl(parcela_pre)} com ac\u00famulo do lance embutido de {emb_str}")
            continue

        # Timeline step 2: "Mes 12"
        if text.startswith("M\u00eas ") and len(text) < 10:
            _replace_full_shape_text(shape, f"M\u00eas {cenario['contemplacao_mes']}")
            continue

        # Lance breakdown
        if "Lance embutido" in text and "R$" in text and "=" in text:
            emb_pct = cenario["lance_embutido_pct"]
            emb_val = cenario["lance_embutido_valor"]
            prop_pct = cenario.get("lance_proprio_pct", 0)
            prop_val = cenario.get("lance_proprio_valor", 0)
            total = cenario.get("total_lances", emb_val + prop_val)
            _replace_full_shape_text(shape,
                f"Lance embutido {fmt_pct_inteiro(emb_pct)} ({fmt_brl(emb_val, com_centavos=False)}) "
                f"+ lance pr\u00f3prio {fmt_pct_inteiro(prop_pct)} ({fmt_brl(prop_val, com_centavos=False)}) "
                f"= {fmt_brl(total, com_centavos=False)}")
            continue

        # Timeline step 3: "Meses 13 a 225"
        if re.match(r"Meses \d+ \u00e0", text):
            inicio = cenario["contemplacao_mes"] + 1
            fim = cenario["prazo_total"]
            _replace_full_shape_text(shape, f"Meses {inicio} \u00e0 {fim}")
            continue

        # Pos-contemplacao descricao
        if text.startswith("Parcelas a partir de R$"):
            idx = cenario["contemplacao_mes"]
            parcela_pos = cenario["parcelas"][idx] if len(cenario["parcelas"]) > idx else 0
            corr_str = fmt_pct_inteiro(cenario["correcao_anual"])
            _replace_full_shape_text(shape, f"Parcelas a partir de {fmt_brl(parcela_pos)} com corre\u00e7\u00e3o anual de {corr_str} pr\u00e9-fixado")
            continue


def _process_proposta_parcelas(slide, all_parcelas: List[float],
                                page_num: int, total_pages: int,
                                correcao: float):
    """Slides de parcelas da Proposta."""
    # Update title
    for shape in slide.shapes:
        text = _get_shape_text(shape)
        if "Tabela de parcelas" in text and "P\u00e1gina" in text:
            _replace_full_shape_text(shape, f"Tabela de parcelas  \u2014  P\u00e1gina {page_num} de {total_pages}")
        elif "Corre\u00e7\u00e3o anual" in text:
            corr_str = fmt_pct_inteiro(correcao)
            _replace_full_shape_text(shape, f"Corre\u00e7\u00e3o anual pr\u00e9-fixada de {corr_str}  |  Inclui seguro")

    # Replace parcela values (the shape numbers match the parcela numbers)
    _replace_parcela_values(slide, all_parcelas)


def _process_contato(slide, data: dict):
    """Ultimo slide: contatos."""
    shapes_ordered = sorted(
        [s for s in slide.shapes if s.has_text_frame and _get_shape_text(s)],
        key=lambda s: (s.top, s.left),
    )

    contact_name_shapes = []
    contact_phone_shapes = []

    found_contatos = False
    for shape in shapes_ordered:
        text = _get_shape_text(shape)
        if "Contatos" in text:
            found_contatos = True
            continue
        if not found_contatos:
            continue
        if "SOMUS" in text or "somuscapital" in text:
            continue

        # Telefone: "(XX) XXXXX-XXXX"
        if re.match(r"\(\d{2}\)\s*\d{4,5}-\d{4}", text.strip()):
            contact_phone_shapes.append(shape)
        # Nome: texto curto sem digitos
        elif (len(text) < 40 and not any(c.isdigit() for c in text)
              and "|" not in text and "Av." not in text and "Rua" not in text
              and "CEP" not in text):
            contact_name_shapes.append(shape)

    if len(contact_name_shapes) >= 1:
        _replace_full_shape_text(contact_name_shapes[0], data.get("contato_1_nome", ""))
    if len(contact_name_shapes) >= 2:
        _replace_full_shape_text(contact_name_shapes[1], data.get("contato_2_nome", ""))
    if len(contact_phone_shapes) >= 1:
        _replace_full_shape_text(contact_phone_shapes[0], data.get("contato_1_telefone", ""))
    if len(contact_phone_shapes) >= 2:
        _replace_full_shape_text(contact_phone_shapes[1], data.get("contato_2_telefone", ""))


# ---------------------------------------------------------------------------
# Funcoes principais de geracao
# ---------------------------------------------------------------------------


def _validate_subtipo(subtipo: str):
    """Valida se o subtipo e valido."""
    if subtipo not in SUBTIPOS_VALIDOS:
        raise ValueError(
            f"Subtipo '{subtipo}' inv\u00e1lido. "
            f"V\u00e1lidos: {', '.join(SUBTIPOS_VALIDOS)}"
        )


def _get_template_path(tipo: str, subtipo: str) -> str:
    """Retorna o caminho do template."""
    filename = f"Demo_{tipo}_{subtipo}.pptx"
    path = os.path.join(TEMPLATES_DIR, filename)
    if not os.path.exists(path):
        raise FileNotFoundError(f"Template n\u00e3o encontrado: {path}")
    return path


def _generate_output_filename(tipo: str, subtipo: str, cliente: str) -> str:
    """Gera o nome do arquivo de saida."""
    cliente_clean = re.sub(r"[^\w\s-]", "", cliente).strip().replace(" ", "_")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"{tipo}_{subtipo}_{cliente_clean}_{timestamp}.pptx"


def gerar_comparativa_pptx(data: dict, subtipo: str, output_dir: str = None) -> str:
    """
    Gera uma apresentacao Comparativa PPTX a partir do template.

    Args:
        data: Dicionario com os dados da comparativa (3 cenarios).
        subtipo: Subtipo do template (CETHD, Lance_Fixo, etc.).
        output_dir: Diretorio de saida. Se None, usa o padrao.

    Returns:
        Caminho completo do arquivo gerado.

    Raises:
        ValueError: Se o subtipo for invalido ou dados inconsistentes.
        FileNotFoundError: Se o template nao existir.
    """
    _validate_subtipo(subtipo)

    if output_dir is None:
        output_dir = DEFAULT_OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)

    template_path = _get_template_path("Comparativa", subtipo)
    output_filename = _generate_output_filename("Comparativa", subtipo, data["cliente_nome"])
    output_path = os.path.join(output_dir, output_filename)

    shutil.copy2(template_path, output_path)

    prs = Presentation(output_path)
    slides = list(prs.slides)

    if len(slides) < 10:
        raise ValueError(f"Template esperava 10 slides, encontrou {len(slides)}.")

    cenarios = data["cenarios"]
    if len(cenarios) != 3:
        raise ValueError(f"Comparativa requer 3 cen\u00e1rios, recebeu {len(cenarios)}.")

    # Slide 1: Capa
    _process_comparativa_cover(slides[0], data)

    # Slide 2: Comparativo das visoes
    _process_comparativa_visoes(slides[1], data)

    # Slides 3-8: Operacoes e parcelas (3 pares)
    for i in range(3):
        _process_comparativa_operacao(slides[2 + i * 2], cenarios[i], i + 1)
        _process_comparativa_parcelas(slides[3 + i * 2], cenarios[i], i + 1)

    # Slide 9: Observacoes (texto padrao - nao modificar)

    # Slide 10: Contato
    _process_contato(slides[9], data)

    prs.save(output_path)
    logger.info(f"Comparativa gerada: {output_path}")
    return output_path


def gerar_proposta_pptx(data: dict, subtipo: str, output_dir: str = None) -> str:
    """
    Gera uma apresentacao Proposta PPTX a partir do template.

    Args:
        data: Dicionario com os dados da proposta (1 cenario).
        subtipo: Subtipo do template (CETHD, PF_Tradicional, etc.).
        output_dir: Diretorio de saida. Se None, usa o padrao.

    Returns:
        Caminho completo do arquivo gerado.

    Raises:
        ValueError: Se o subtipo for invalido ou dados inconsistentes.
        FileNotFoundError: Se o template nao existir.
    """
    _validate_subtipo(subtipo)

    if output_dir is None:
        output_dir = DEFAULT_OUTPUT_DIR
    os.makedirs(output_dir, exist_ok=True)

    template_path = _get_template_path("Proposta", subtipo)
    output_filename = _generate_output_filename("Proposta", subtipo, data["cliente_nome"])
    output_path = os.path.join(output_dir, output_filename)

    shutil.copy2(template_path, output_path)

    prs = Presentation(output_path)
    slides = list(prs.slides)

    tem_contexto = subtipo not in SUBTIPOS_SEM_CONTEXTO
    min_slides = 11 if tem_contexto else 10

    if len(slides) < min_slides:
        raise ValueError(f"Template esperava {min_slides} slides, encontrou {len(slides)}.")

    cenario = data["cenario"]
    offset = 0 if tem_contexto else -1

    # Slide 1: Capa
    _process_proposta_cover(slides[0], data)

    # Slide 2: Contexto (se existir)
    if tem_contexto:
        _process_proposta_contexto(slides[1], data)

    # Slide 3 (ou 2): Principais dados
    _process_proposta_dados(slides[2 + offset], data)

    # Slide 4 (ou 3): Estrutura
    _process_proposta_estrutura(slides[3 + offset], data)

    # Slides de parcelas
    parcelas_start = 4 + offset
    parcelas = cenario["parcelas"]
    correcao = cenario["correcao_anual"]

    # Contar slides de parcelas
    num_parcela_slides = 0
    for si in range(parcelas_start, len(slides)):
        slide_text = ""
        for shape in slides[si].shapes:
            if shape.has_text_frame:
                slide_text += shape.text_frame.text
        if "Tabela de parcelas" in slide_text or "Parcela 1" in slide_text or "Parcela" in slide_text:
            # Verificar que e realmente um slide de parcelas e nao observacoes
            if "Observa\u00e7\u00f5es" not in slide_text and "Fale com" not in slide_text:
                num_parcela_slides += 1
            else:
                break
        else:
            break

    total_pages = num_parcela_slides
    for page in range(num_parcela_slides):
        _process_proposta_parcelas(
            slides[parcelas_start + page], parcelas,
            page + 1, total_pages, correcao
        )

    # Ultimo slide: Contato
    _process_contato(slides[-1], data)

    prs.save(output_path)
    logger.info(f"Proposta gerada: {output_path}")
    return output_path


# ---------------------------------------------------------------------------
# Bloco de teste
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import sys

    logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

    print("=" * 60)
    print("Teste do gerador de PPTX - Somus Capital")
    print("=" * 60)

    # -----------------------------------------------------------------------
    # Dados de exemplo - COMPARATIVA
    # -----------------------------------------------------------------------
    parcelas_48_c1 = (
        [12_300.00] * 6 + [7_100.00] * 18 + [7_313.00] * 12 + [7_532.39] * 12
    )
    parcelas_48_c2 = (
        [10_850.00] * 12 + [6_400.00] * 12 + [6_592.00] * 12 + [6_789.76] * 12
    )
    parcelas_48_c3 = (
        [9_200.00] * 24 + [5_850.00] * 12 + [6_025.50] * 12
    )

    data_comparativa = {
        "cliente_nome": "Carlos Mendes",
        "data_mes_ano": "Mar\u00e7o 2026",
        "cenarios": [
            {
                "contemplacao_mes": 6,
                "cet_mensal": 0.6210,
                "cet_anual": 7.71,
                "lance_livre_pct": 25,
                "lance_livre_valor": 375_000,
                "credito_disponivel": 1_050_000,
                "alavancagem": 675_000,
                "valor_carta": 1_500_000,
                "lance_embutido_pct": 30,
                "lance_embutido_valor": 450_000,
                "taxa_adm": 19.00,
                "fundo_reserva": 3.00,
                "prazo_total": 48,
                "correcao_anual": 3,
                "parcelas": parcelas_48_c1,
            },
            {
                "contemplacao_mes": 12,
                "cet_mensal": 0.6854,
                "cet_anual": 8.52,
                "lance_livre_pct": 15,
                "lance_livre_valor": 225_000,
                "credito_disponivel": 1_050_000,
                "alavancagem": 825_000,
                "valor_carta": 1_500_000,
                "lance_embutido_pct": 30,
                "lance_embutido_valor": 450_000,
                "taxa_adm": 19.00,
                "fundo_reserva": 3.00,
                "prazo_total": 48,
                "correcao_anual": 3,
                "parcelas": parcelas_48_c2,
            },
            {
                "contemplacao_mes": 24,
                "cet_mensal": 0.7512,
                "cet_anual": 9.38,
                "lance_livre_pct": 5,
                "lance_livre_valor": 75_000,
                "credito_disponivel": 1_050_000,
                "alavancagem": 975_000,
                "valor_carta": 1_500_000,
                "lance_embutido_pct": 30,
                "lance_embutido_valor": 450_000,
                "taxa_adm": 19.00,
                "fundo_reserva": 3.00,
                "prazo_total": 48,
                "correcao_anual": 3,
                "parcelas": parcelas_48_c3,
            },
        ],
        "contato_1_nome": "Marcos Silva",
        "contato_1_telefone": "(21) 99876-5432",
        "contato_2_nome": "Julia Santos",
        "contato_2_telefone": "(11) 98765-4321",
    }

    # -----------------------------------------------------------------------
    # Dados de exemplo - PROPOSTA
    # -----------------------------------------------------------------------
    parcelas_225 = []
    valor_base = 16_554.16
    valor_pos = 8_662.20
    for m in range(1, 226):
        if m <= 12:
            parcelas_225.append(valor_base)
        else:
            anos = (m - 13) // 12
            parcelas_225.append(round(valor_pos * (1.03 ** anos), 2))

    data_proposta = {
        "cliente_nome": "Ricardo Almeida",
        "data_mes_ano": "Mar\u00e7o 2026",
        "empresa_nome": "Almeida Holdings",
        "localizacao": "Rio de Janeiro, RJ",
        "descricao_operacao": (
            "Opera\u00e7\u00e3o de cons\u00f3rcio imobili\u00e1rio estruturada com carta de cr\u00e9dito "
            "de R$ 4.326.000, administrada pelo Ita\u00fa. A estrat\u00e9gia combina lance "
            "embutido de 30% para contempla\u00e7\u00e3o acelerada no m\u00eas 12, com corre\u00e7\u00e3o "
            "anual pr\u00e9-fixada de 3% ao longo dos 225 meses de opera\u00e7\u00e3o."
        ),
        "objetivo_titulo": "Aquisi\u00e7\u00e3o de im\u00f3vel comercial via cons\u00f3rcio",
        "horizonte_texto": "Horizonte de 225 meses (18 anos e 9 meses)",
        "objetivo_descricao": (
            "Capta\u00e7\u00e3o de cr\u00e9dito de R$ 4.326.000 via cons\u00f3rcio estruturado, "
            "preservando caixa e otimizando o custo de aquisi\u00e7\u00e3o."
        ),
        "administradora": "Ita\u00fa",
        "cenario": {
            "contemplacao_mes": 12,
            "cet_mensal": 0.7812,
            "cet_anual": 9.79,
            "lance_proprio_pct": 47,
            "lance_proprio_valor": 2_033_220,
            "lance_embutido_pct": 30,
            "lance_embutido_valor": 1_297_800,
            "credito_disponivel": 3_028_200,
            "credito_efetivo": 994_980,
            "valor_carta": 4_326_000,
            "taxa_adm": 20.00,
            "fundo_reserva": 3.00,
            "prazo_total": 225,
            "correcao_anual": 3,
            "parcelas": parcelas_225,
            "estrutura_descricao": (
                "A opera\u00e7\u00e3o utiliza cons\u00f3rcio imobili\u00e1rio com carta de cr\u00e9dito "
                "de R$ 4.326.000, administrado pelo Ita\u00fa. A estrat\u00e9gia combina lance "
                "embutido de 30% (R$ 1.297.800) com lance pr\u00f3prio de 47% (R$ 2.033.220), "
                "totalizando R$ 3.331.020 em lances."
            ),
            "cronograma_resumo": (
                "Contempla\u00e7\u00e3o no m\u00eas 12 via lance embutido (30%) + lance pr\u00f3prio (47%), "
                "seguida de amortiza\u00e7\u00e3o em 213 meses com parcelas corrigidas anualmente "
                "em 3% pr\u00e9-fixado."
            ),
            "captacao_meses": 12,
            "total_lances": 3_331_020,
        },
        "contato_1_nome": "Marcos Silva",
        "contato_1_telefone": "(21) 99876-5432",
        "contato_2_nome": "Julia Santos",
        "contato_2_telefone": "(11) 98765-4321",
    }

    # -----------------------------------------------------------------------
    # Gerar arquivos
    # -----------------------------------------------------------------------
    try:
        print("\n[1/2] Gerando Comparativa CETHD...")
        path_comp = gerar_comparativa_pptx(data_comparativa, "CETHD")
        print(f"  -> {path_comp}")
    except Exception as e:
        print(f"  -> ERRO: {e}")
        import traceback
        traceback.print_exc()

    try:
        print("\n[2/2] Gerando Proposta CETHD...")
        path_prop = gerar_proposta_pptx(data_proposta, "CETHD")
        print(f"  -> {path_prop}")
    except Exception as e:
        print(f"  -> ERRO: {e}")
        import traceback
        traceback.print_exc()

    print("\n" + "=" * 60)
    print("Teste concluido.")
    print("=" * 60)
